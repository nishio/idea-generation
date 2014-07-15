#!/usr/bin/env python
# -*- encoding: utf-8 -*-
"""
テキストをpptxにする
改行位置制御や図を貼りたいなどの細かいニーズに答えるのは
ユーザの手元のPowerPointでやってもらおうという発想
"""
import argparse

from pptx import Presentation
from pptx.util import Inches, Pt, Cm
from pptx.enum.text import MSO_AUTO_SIZE, PP_ALIGN, MSO_ANCHOR

def main():
    parser = argparse.ArgumentParser(description='Text to PPTX.')
    parser.add_argument('--test', action='store_true')

    args = parser.parse_args()
    if args.test:
        texts = ['あ' * x for x in range(1, 101)]
    else:
        import sys
        texts = sys.stdin.read().split('\n')

    make_pptx(texts)

def find_best_fontsize(text):
    sizes = [415, 415, 346, 240]
    chars = len(text.decode('utf-8').encode('sjis')) / 2
    if chars < len(sizes):
        return sizes[chars]

    # means 'if chars leq 6 (2 * 3), fontsize is 200pt'
    sizes = [
        (2 * 3, 200), (2 * 4, 167), (3 * 4, 159),
        (3 * 5, 125), (4 * 6, 110), (5 * 8, 90),
        (5 * 9, 80), (6 * 10, 70), (7 * 14, 60),
    ]

    for lim, siz in sizes:
        if chars <= lim:
            return siz

    return 50


def make_pptx(texts):
    prs = Presentation()
    for text in texts:
        blank_slide_layout = prs.slide_layouts[6]
        slide = prs.slides.add_slide(blank_slide_layout)

        txBox = slide.shapes.add_textbox(0, 0, Cm(25.4), Cm(19.1))
        tf = txBox.textframe
        tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
        tf.word_wrap = True
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = tf.paragraphs[0]
        p.text = text
        p.font.size = Pt(find_best_fontsize(text))

        p.alignment = PP_ALIGN.CENTER
    prs.save('test.pptx')


if __name__ == '__main__':
    main()

